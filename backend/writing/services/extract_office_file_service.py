import os
import re
import sys
from datetime import datetime
from io import BytesIO
from typing import Any

# For pdf
import fitz
from fastapi import HTTPException
from llama_index.core import Document

# For pptx, ppt
from pptx import Presentation as PPT_Presentation

# For docx, doc
from spire.doc import Document as DocDocument
from spire.doc import DocumentObjectType
from spire.doc.common import Stream as DocStream
from writing.services.base_service import TextExtractor
from writing.helpers.convert_ppt_helper import convert_ppt_to_pptx_stream
from writing.helpers.pdf_helper import get_tessocr, is_scan_page, is_two_column_paper
from writing.helpers.word_table_helper import extract_text_from_table
from writing.utils.format_time_utils import _format_file_timestamp


class OfficeFile(TextExtractor):
    def __init__(self):
        self.invalide_unicode = chr(0xFFFD)

    async def extract_text(self, file):
        file_content = await file.read()
        file_name = file.filename
        _, file_type = os.path.splitext(file_name)
        if file_type in [".txt"]:
            texts = self.extract_text_from_txt(file_content)
            return texts
        elif file_type in [".doc", ".docx"]:
            texts = self.extract_text_from_docx(file_content)
            return texts
        elif file_type == ".pdf":
            texts = self.extract_text_from_pdf(file_content)
            return texts
        elif file_type == ".pptx":
            texts = self.extract_text_from_powerpoint(file_content)
            return texts
        elif file_type == ".ppt":
            file_content = convert_ppt_to_pptx_stream(file)
            texts = self.extract_text_from_powerpoint(file_content)
            return texts

    def supports_file_type(self, file_name):
        _, file_extension = os.path.splitext(file_name)
        return file_extension in [".txt", ".doc", ".docx", ".pdf", ".pptx", ".ppt"]

    def create_docs(self, texts: list[Any], file_name: str, user_id: str):
        docs = []
        for i, text in enumerate(texts):
            metadata = {}
            metadata["total_pages"] = len(texts)
            metadata["creation_date"] = _format_file_timestamp(
                timestamp=datetime.now().timestamp(), include_time=True
            )
            metadata["file_name"] = str(file_name)
            metadata["user_id"] = user_id
            metadata["source"] = i + 1
            doc = Document(text=text, extra_info=metadata)
            docs.append(doc)
        return docs

    def extract_text_from_pdf(self, document_content):
        try:
            document_stream = BytesIO(document_content)
            doc = fitz.open(stream=document_stream, filetype="pdf")
            is_two_column, list_res = is_two_column_paper(doc)
            list_page_text = []
            for page_number, page in enumerate(doc):
                if is_scan_page(page):  # For Scan PDF
                    bbox = page.rect
                    new_text = get_tessocr(page, bbox)
                    list_page_text.append(new_text)
                else:  # For normal PDF
                    page_text = ""
                    res = list_res[page_number]
                    combined = list(zip(res["bbox"], res["text"], strict=False))
                    if not is_two_column:
                        sorted_combined = sorted(
                            combined, key=lambda x: (x[0][1], x[0][0])
                        )
                        res["bbox"], res["text"] = zip(*sorted_combined, strict=False)
                        res["bbox"] = list(res["bbox"])
                        res["text"] = list(res["text"])
                    for bbox, text in zip(res["bbox"], res["text"], strict=False):
                        if self.invalide_unicode in text:
                            text1 = text.lstrip()
                            sb = " " * (len(text) - len(text1))  # leading spaces
                            text1 = text.rstrip()
                            sa = " " * (len(text) - len(text1))  # trailing spaces
                            new_text = sb + get_tessocr(page, bbox) + sa
                            page_text += new_text + "\n"
                        else:
                            page_text += text + "\n"
                    list_page_text.append(page_text)
            return list_page_text
        except HTTPException as e:
            raise e

    def extract_text_from_docx(self, document_content):
        try:
            docstream = DocStream(document_content)
            doc = DocDocument(docstream)
            list_section_content = []
            for s in range(doc.Sections.Count):
                elements = doc.Sections[s].Body.ChildObjects
                section_content = ""
                for i in range(elements.Count):
                    element = elements.get_Item(i)
                    if element.DocumentObjectType == DocumentObjectType.Paragraph:
                        paragraph = element
                        section_content += paragraph.Text + "\n"

                    elif element.DocumentObjectType == DocumentObjectType.Table:
                        table = element
                        table_text = extract_text_from_table(table)
                        section_content += table_text + "\n"
                section_content = re.sub(r"\x0b", "\n", section_content)
                list_section_content.append(section_content)
            return list_section_content
        except HTTPException as e:
            raise e

    def extract_text_from_powerpoint(self, document_content):
        try:
            document_stream = BytesIO(document_content)
            presentation = PPT_Presentation(document_stream)
            list_page_text = []
            for slide in presentation.slides:
                page_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        page_text += shape.text + "\n"
                list_page_text.append(page_text)
            return list_page_text
        except HTTPException as e:
            raise e

    def extract_text_from_txt(self, document_content):
        try:
            file_like_object = BytesIO(document_content)
            text = file_like_object.read().decode("utf-8")
            return [text]
        except HTTPException as e:
            raise e
